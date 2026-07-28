"""PPTX reader (optional: requires ``python-pptx``).

The optional dependency is imported lazily so the package imports cleanly even
when ``python-pptx`` is not installed; :meth:`read` raises a clear ``RuntimeError``
explaining how to enable ``.pptx`` support. Each slide becomes one "page".
"""

from __future__ import annotations

from agents.readers.base import BaseReader, ReaderResult

try:
    from pptx import Presentation as _PptxPresentation  # type: ignore

    _PPTX_AVAILABLE = True
except Exception:  # noqa: BLE001 - optional dependency, import is best-effort
    _PPTX_AVAILABLE = False


class PptxReader(BaseReader):
    """Extracts text from each slide of a ``.pptx`` presentation."""

    extensions = (".pptx",)

    def read(self, data: bytes, filename: str = "") -> ReaderResult:
        if not _PPTX_AVAILABLE:
            raise RuntimeError(
                "PptxReader requires the optional 'python-pptx' package. "
                "Install it with: pip install python-pptx"
            )

        import io

        presentation = _PptxPresentation(io.BytesIO(data))
        slide_texts: list[str] = []
        for slide in presentation.slides:
            parts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = para.text.strip()
                        if line:
                            parts.append(line)
            if parts:
                slide_texts.append("\n".join(parts))

        raw_text = "\n\n".join(slide_texts)
        return ReaderResult(
            raw_text=raw_text,
            pages=slide_texts,
            num_pages=len(slide_texts),
            format="pptx",
        )
