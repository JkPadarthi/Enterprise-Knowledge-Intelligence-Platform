"""Tests for the pluggable document readers (multi-format ingestion).

Markdown and HTML backends use only the standard library, so they run in any
environment. DOCX/PPTX backends require optional packages and are skipped when
those are not installed.
"""
from __future__ import annotations

import io

import fitz
import pytest

from agents.reader import ReaderAgent
from agents.readers import (
    HTMLReader,
    MarkdownReader,
    PDFReader,
    ReaderRegistry,
    get_reader,
    supported_extensions,
)
from agents.readers.docx import DocxReader, _DOCX_AVAILABLE
from agents.readers.pptx import PptxReader, _PPTX_AVAILABLE
from models.schema import AgentState


def _make_pdf(text: str) -> bytes:
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((50, 100), text)
    data = doc.tobytes()
    doc.close()
    return data


# ── Markdown ───────────────────────────────────────────────────────────────
def test_markdown_reader_returns_text():
    reader = MarkdownReader()
    res = reader.read("# Title\n\nSome **content**.".encode("utf-8"), "doc.md")
    assert res.format == "markdown"
    assert "Title" in res.raw_text
    assert res.num_pages == 1
    assert res.pages == [res.raw_text]


def test_markdown_reader_handles_txt():
    reader = MarkdownReader()
    res = reader.read(b"plain text file", "notes.txt")
    assert res.raw_text == "plain text file"
    assert ".txt" in MarkdownReader.extensions


# ── HTML ───────────────────────────────────────────────────────────────────
def test_html_reader_strips_tags():
    html = b"<html><head><style>.a{}</style></head><body><h1>Hi</h1><p>Hello <b>world</b></p><script>bad()</script></body></html>"
    res = HTMLReader().read(html, "page.html")
    assert res.format == "html"
    assert "Hi" in res.raw_text
    assert "Hello" in res.raw_text
    assert "world" in res.raw_text
    assert "<" not in res.raw_text  # no tags survive
    assert "bad()" not in res.raw_text  # script content dropped


def test_html_reader_decodes_entities():
    res = HTMLReader().read(b"<p>caf&eacute; &amp; co</p>", "x.html")
    assert "café" in res.raw_text
    assert "&" in res.raw_text  # the literal ampersand from &amp;


# ── PDF ──────────────────────────────────────────────────────────────────
def test_pdf_reader_extracts_text():
    res = PDFReader().read(_make_pdf("Hello PDF"), "doc.pdf")
    assert res.format == "pdf"
    assert res.num_pages == 1
    assert "Hello PDF" in res.raw_text


def test_pdf_reader_rejects_garbage():
    with pytest.raises(ValueError):
        PDFReader().read(b"not a pdf", "doc.pdf")


# ── Registry / factory ──────────────────────────────────────────────────────
def test_get_reader_selects_by_extension():
    assert isinstance(get_reader("a.pdf"), PDFReader)
    assert isinstance(get_reader("a.MD"), MarkdownReader)
    assert isinstance(get_reader("a.html"), HTMLReader)
    assert isinstance(get_reader("a.docx"), DocxReader)
    assert isinstance(get_reader("a.pptx"), PptxReader)


def test_get_reader_defaults_to_pdf():
    # No extension or unknown extension falls back to the PDF reader.
    assert isinstance(get_reader(""), PDFReader)
    assert isinstance(get_reader("weird.xyz"), PDFReader)


def test_supported_extensions_non_empty():
    exts = supported_extensions()
    assert ".pdf" in exts
    assert ".md" in exts
    assert ".html" in exts
    assert ".docx" in exts
    assert ".pptx" in exts


def test_registry_is_extensible():
    class StubReader(MarkdownReader):
        extensions = (".stub",)

    reg = ReaderRegistry()
    reg.register(StubReader())
    assert isinstance(reg.get(".stub"), StubReader)


# ── Optional DOCX / PPTX ─────────────────────────────────────────────────────
@pytest.mark.skipif(_DOCX_AVAILABLE, reason="python-docx is installed")
def test_docx_reader_reports_missing_dependency():
    with pytest.raises(RuntimeError, match="python-docx"):
        DocxReader().read(b"", "doc.docx")


@pytest.mark.skipif(_PPTX_AVAILABLE, reason="python-pptx is installed")
def test_pptx_reader_reports_missing_dependency():
    with pytest.raises(RuntimeError, match="python-pptx"):
        PptxReader().read(b"", "doc.pptx")


@pytest.mark.skipif(not _DOCX_AVAILABLE, reason="python-docx not installed")
def test_docx_reader_reads_paragraphs():
    from docx import Document

    doc = Document()
    doc.add_paragraph("First paragraph.")
    doc.add_paragraph("Second paragraph.")
    bio = io.BytesIO()
    doc.save(bio)
    res = DocxReader().read(bio.getvalue(), "doc.docx")
    assert "First paragraph." in res.raw_text
    assert "Second paragraph." in res.raw_text


@pytest.mark.skipif(not _PPTX_AVAILABLE, reason="python-pptx not installed")
def test_pptx_reader_reads_slides():
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    box = slide.shapes.add_textbox(0, 0, 100, 100)
    box.text = "Slide One"
    bio = io.BytesIO()
    prs.save(bio)
    res = PptxReader().read(bio.getvalue(), "deck.pptx")
    assert "Slide One" in res.raw_text
    assert res.num_pages >= 1


# ── ReaderAgent routing ──────────────────────────────────────────────────────
@pytest.mark.asyncio
async def test_reader_agent_routes_markdown():
    agent = ReaderAgent()
    state = AgentState(file_bytes=b"# Heading\n\nbody text", filename="note.md")
    result = await agent.run(state)
    assert "Heading" in result["raw_text"]
    assert result["num_pages"] == 1


@pytest.mark.asyncio
async def test_reader_agent_pdf_still_works():
    agent = ReaderAgent()
    state = AgentState(file_bytes=_make_pdf("Legacy PDF path"), filename="doc.pdf")
    result = await agent.run(state)
    assert "Legacy PDF path" in result["raw_text"]
